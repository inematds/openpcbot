#!/usr/bin/env python3
"""
Process Documents -> Obsidian Markdown Notes
Fallback order: Gemini (n14) -> Claude CLI -> Ollama local

Usage:
    python scripts/process_docs_to_obsidian.py <input_folder> <output_folder>

Example:
    python scripts/process_docs_to_obsidian.py ~/Documents/files ~/vault/inbox

Supported: PDF, PPTX, DOCX, TXT, MD
"""

import os
import sys
import json
import subprocess
from pathlib import Path
from datetime import date

from dotenv import load_dotenv
load_dotenv()

GEMINI_MODEL = "gemini-2.0-flash"
API_KEY = os.environ.get("GOOGLE_API_KEY")
OLLAMA_URL = os.environ.get("OLLAMA_URL", "http://localhost:11434")
OLLAMA_MODEL = os.environ.get("OLLAMA_MODEL", "qwen3.5:35b-a3b")
TODAY = date.today().isoformat()

SUPPORTED = {".pdf", ".pptx", ".ppt", ".docx", ".doc", ".txt", ".md"}

SYNTHESIS_PROMPT = """
You are processing a document to be stored in an Obsidian second brain vault.

Your job is to extract SIGNAL and discard NOISE.

SIGNAL = key insights, decisions, frameworks, facts, action items, names, dates that matter
NOISE  = headers/footers, legal boilerplate, filler sentences, repeated content, formatting artifacts

Output ONLY a clean Obsidian-compatible Markdown note with this exact structure:

---
type: [note | meeting | report | presentation | research | reference | other]
topic: [2-4 word topic description]
source: {filename}
date_processed: {date}
tags: [comma-separated lowercase tags, no #]
---

# [Concise, descriptive title — what IS this document?]

## Key Insights
[3-7 bullet points of the highest-signal takeaways]

## Context
[1-2 sentences: what is this document, who created it, what was it for?]

## Details Worth Keeping
[Any specific data, frameworks, quotes, or facts that should be preserved verbatim]

## Action Items
[Only include if action items exist. Delete this section if none.]

---

RULES:
- Total note length: 300-600 words MAX. Ruthlessly compress.
- Do not copy paste large blocks of text. Synthesize.
- If the document is very short (under 200 words), keep it mostly intact.
- Write in clean, clear prose. No corporate speak.
- If you cannot read the document or it's empty, return: "# [filename] — Could not process"
"""


# ── File readers ─────────────────────────────────────────────────────────────

def read_pdf_text(file_path: Path) -> str:
    try:
        import pdfplumber
        with pdfplumber.open(file_path) as pdf:
            return "\n\n".join(p.extract_text() or "" for p in pdf.pages).strip()
    except ImportError:
        return f"[PDF — {file_path.stat().st_size:,} bytes. Install pdfplumber.]"
    except Exception as e:
        return f"[PDF error: {e}]"


def read_pptx_text(file_path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = [shape.text.strip() for shape in slide.shapes
                     if hasattr(shape, "text") and shape.text.strip()]
            if texts:
                slides.append(f"[Slide {i}] " + " | ".join(texts))
        return "\n".join(slides)
    except ImportError:
        return "[Error: install python-pptx]"


def read_docx_text(file_path: Path) -> str:
    try:
        import docx
        doc = docx.Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except ImportError:
        return "[Error: install python-docx]"


def read_text(file_path: Path) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


# ── LLM backends (fallback: Gemini -> Claude -> Ollama) ──────────────────────

def _call_gemini(prompt: str, pdf_bytes: bytes | None = None) -> str | None:
    if not API_KEY:
        return None
    try:
        from google import genai
        from google.genai import types
        client = genai.Client(api_key=API_KEY)
        if pdf_bytes:
            contents = [
                types.Part.from_bytes(data=pdf_bytes, mime_type="application/pdf"),
                prompt
            ]
        else:
            contents = [prompt]
        response = client.models.generate_content(model=GEMINI_MODEL, contents=contents)
        return response.text.strip()
    except Exception as e:
        print(f"  [Gemini failed: {e}]")
        return None


def _call_claude(prompt: str) -> str | None:
    try:
        result = subprocess.run(
            ["claude", "-p", prompt],
            capture_output=True, text=True, timeout=120
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout.strip()
        return None
    except (FileNotFoundError, subprocess.TimeoutExpired) as e:
        print(f"  [Claude CLI failed: {e}]")
        return None


def _call_ollama(prompt: str) -> str | None:
    import urllib.request
    try:
        payload = json.dumps({
            "model": OLLAMA_MODEL,
            "prompt": prompt,
            "stream": False,
        }).encode()
        req = urllib.request.Request(
            f"{OLLAMA_URL}/api/generate",
            data=payload,
            headers={"Content-Type": "application/json"},
        )
        with urllib.request.urlopen(req, timeout=120) as resp:
            data = json.loads(resp.read())
            return data.get("response", "").strip() or None
    except Exception as e:
        print(f"  [Ollama failed: {e}]")
        return None


def call_llm(prompt: str, pdf_bytes: bytes | None = None) -> tuple[str, str]:
    """Returns (result, backend_name)."""
    result = _call_gemini(prompt, pdf_bytes)
    if result:
        return result, "Gemini"

    print("  Gemini unavailable, trying Claude CLI...")
    result = _call_claude(prompt)
    if result:
        return result, "Claude"

    print("  Claude unavailable, trying Ollama...")
    result = _call_ollama(prompt)
    if result:
        return result, "Ollama"

    return "[All LLM backends failed.]", "none"


# ── Main processing ──────────────────────────────────────────────────────────

def process_file(file_path: Path) -> tuple[str | None, str]:
    suffix = file_path.suffix.lower()
    prompt = SYNTHESIS_PROMPT.format(filename=file_path.name, date=TODAY)

    if suffix == ".pdf":
        with open(file_path, "rb") as f:
            pdf_bytes = f.read()
        text = read_pdf_text(file_path)
        # Try Gemini with native PDF first, fallback to text extraction
        result, backend = call_llm(prompt, pdf_bytes=pdf_bytes)
        if backend == "Gemini":
            return result, backend
        # For Claude/Ollama, use extracted text
        text_prompt = f"DOCUMENT CONTENT:\n\n{text}\n\n{prompt}"
        return call_llm(text_prompt)
    elif suffix in {".pptx", ".ppt"}:
        text = read_pptx_text(file_path)
        return call_llm(f"PRESENTATION CONTENT:\n\n{text}\n\n{prompt}")
    elif suffix in {".docx", ".doc"}:
        text = read_docx_text(file_path)
        return call_llm(f"DOCUMENT CONTENT:\n\n{text}\n\n{prompt}")
    elif suffix in {".txt", ".md"}:
        text = read_text(file_path)
        return call_llm(f"TEXT CONTENT:\n\n{text}\n\n{prompt}")
    else:
        return None, "none"


def process_folder(input_folder: str, output_folder: str):
    input_path = Path(input_folder).expanduser()
    output_path = Path(output_folder).expanduser()
    output_path.mkdir(parents=True, exist_ok=True)

    files = [f for f in input_path.iterdir() if f.suffix.lower() in SUPPORTED]

    if not files:
        print(f"No supported files found in {input_path}")
        print(f"Supported types: {', '.join(SUPPORTED)}")
        return

    print(f"\nInput:  {input_path}")
    print(f"Output: {output_path}")
    print(f"Files:  {len(files)} found\n")
    print("-" * 50)

    success, skipped = 0, 0

    for i, file_path in enumerate(files, 1):
        print(f"[{i}/{len(files)}] {file_path.name}")

        result, backend = process_file(file_path)

        if result:
            out_file = output_path / (file_path.stem + ".md")
            with open(out_file, "w", encoding="utf-8") as f:
                f.write(result)
            print(f"  -> {out_file.name} (via {backend})")
            success += 1
        else:
            print(f"  Skipped")
            skipped += 1

    print("-" * 50)
    print(f"\nDone: {success} notes created, {skipped} skipped")
    print(f"Notes in: {output_path}")


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print(__doc__)
        sys.exit(1)
    process_folder(sys.argv[1], sys.argv[2])
