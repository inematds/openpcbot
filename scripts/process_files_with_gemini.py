#!/usr/bin/env python3
"""
Process any folder of files — extract content, summarise, analyse.
Fallback order: Gemini (n14) -> Claude CLI -> Ollama local

Usage:
  python scripts/process_files_with_gemini.py                    # processes inbox/
  python scripts/process_files_with_gemini.py path/to/folder     # processes custom folder

Supported: PDF, PPTX, XLSX, DOCX, CSV, JSON, XML, MD, TXT, PY, JS, HTML, CSS, any text file

Output: outputs/file_summaries/YYYY-MM-DD/
  - <filename>_summary.md    (per-file analysis)
  - MASTER_SUMMARY.md        (full digest of all files)
"""

import os
import sys
import json
import csv
import subprocess
import io
from pathlib import Path
from datetime import date
from dotenv import load_dotenv

load_dotenv()

API_KEY = os.environ.get("GOOGLE_API_KEY")
OLLAMA_URL = os.environ.get("OLLAMA_URL", "http://localhost:11434")
OLLAMA_MODEL = os.environ.get("OLLAMA_MODEL", "qwen3.5:35b-a3b")
GEMINI_MODEL = "gemini-2.0-flash"

BASE_DIR = Path(__file__).parent.parent
TODAY = date.today().isoformat()
OUTPUT_DIR = BASE_DIR / "outputs" / "file_summaries" / TODAY
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# Track which backend was used
_backend_used = None


def extract_pdf(path: Path) -> str:
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            text = "\n\n".join(p.extract_text() or "" for p in pdf.pages)
        return text.strip()
    except ImportError:
        return f"[PDF — {path.stat().st_size:,} bytes. Install pdfplumber for text extraction.]"
    except Exception as e:
        return f"[PDF extraction error: {e}]"


def extract_docx(path: Path) -> str:
    try:
        from docx import Document
        doc = Document(path)
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                style = para.style.name
                if "Heading" in style:
                    parts.append(f"\n## {para.text}")
                else:
                    parts.append(para.text)
        return "\n".join(parts)
    except Exception as e:
        return f"[DOCX extraction error: {e}]"


def extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slides.append(f"[Slide {i}]\n" + "\n".join(texts))
        return "\n\n".join(slides)
    except Exception as e:
        return f"[PPTX extraction error: {e}]"


def extract_xlsx(path: Path) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"## Sheet: {sheet_name}")
            rows = []
            for row in ws.iter_rows(values_only=True):
                if any(c is not None for c in row):
                    rows.append(" | ".join(str(c) if c is not None else "" for c in row))
            parts.extend(rows[:50])
            if ws.max_row > 50:
                parts.append(f"... ({ws.max_row - 50} more rows)")
        return "\n".join(parts)
    except Exception as e:
        return f"[XLSX extraction error: {e}]"


def extract_csv(path: Path) -> str:
    try:
        with open(path, newline="", encoding="utf-8") as f:
            reader = csv.reader(f)
            rows = list(reader)
        lines = [" | ".join(row) for row in rows[:30]]
        if len(rows) > 30:
            lines.append(f"... ({len(rows) - 30} more rows)")
        return "\n".join(lines)
    except Exception as e:
        return f"[CSV extraction error: {e}]"


def extract_text(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8", errors="replace")
    except Exception as e:
        return f"[Text extraction error: {e}]"


EXTRACTORS = {
    ".pdf":  extract_pdf,
    ".docx": extract_docx,
    ".pptx": extract_pptx,
    ".xlsx": extract_xlsx,
    ".xls":  extract_xlsx,
    ".csv":  extract_csv,
}

TEXT_EXTENSIONS = {".txt", ".md", ".json", ".xml", ".py", ".js", ".ts",
                   ".html", ".htm", ".css", ".yaml", ".yml", ".toml", ".sh"}


def extract_content(path: Path) -> tuple[str, str]:
    ext = path.suffix.lower()
    if ext in EXTRACTORS:
        label = ext.lstrip(".").upper()
        return EXTRACTORS[ext](path), label
    elif ext in TEXT_EXTENSIONS:
        return extract_text(path), ext.lstrip(".").upper()
    else:
        try:
            content = path.read_text(encoding="utf-8", errors="replace")
            return content, "TEXT"
        except Exception:
            return f"[Binary file — cannot extract text from {ext}]", "BINARY"


ANALYSIS_PROMPT = """You are writing a quick-reference note for an Obsidian knowledge base.
The note will be scanned — not read. It needs to be instantly useful, conversational, and memorable.

File: {filename}
Type: {file_type}

Extracted content:
---
{content}
---

Silently classify the file as one of these types (do NOT output the classification — just use the correct format):

A) DELIVERABLE — invoice, report, contract, meeting notes, budget, presentation, sales data, proposal
B) REFERENCE — code, config, HTML, CSS, script, template, readme, design system, transcript, guide

Write the note using the matching format below. Output ONLY the note — no preamble, no classification label.

---

FORMAT A — DELIVERABLE FILES:

## TL;DR
One punchy sentence. What is this and what's the key number/date/outcome?

## Numbers & Dates
Only the figures that actually matter. No fluff.
- [key figure or deadline]
- [key figure or deadline]
- [key figure or deadline]

## What This Means
2-3 short, conversational sentences. Not bullet points — write like you're telling a colleague over Slack.
Focus on implications, not just facts. What's surprising, important, or worth flagging?

## Next Move
1-2 specific, non-obvious actions. Skip anything obvious like "open the file" or "read the document."
If nothing meaningful needs doing, omit this section entirely.

---

FORMAT B — REFERENCE FILES:

## TL;DR
One punchy sentence. What does this file do or contain?

## What's Inside
3-5 bullets. Be specific — name actual functions, config keys, sections, or design tokens. No vague descriptions.

## Worth Knowing
1-2 sentences on something non-obvious about how this file works or connects to other things.
If there's nothing worth flagging, omit this section.

---

Rules:
- Write conversationally. Short sentences. No consultant-speak.
- Never use the phrase "this file" more than once.
- No padding. If a section would be filler, leave it out.
- Total response under 250 words."""


# ── LLM backends (fallback: Gemini -> Claude -> Ollama) ──────────────────────

def _call_gemini(prompt: str) -> str | None:
    """Try Gemini API. Returns None on failure."""
    if not API_KEY:
        return None
    try:
        from google import genai
        client = genai.Client(api_key=API_KEY)
        response = client.models.generate_content(
            model=GEMINI_MODEL,
            contents=prompt,
        )
        return response.text.strip()
    except Exception as e:
        print(f"  [Gemini failed: {e}]")
        return None


def _call_claude(prompt: str) -> str | None:
    """Try Claude CLI. Returns None on failure."""
    try:
        result = subprocess.run(
            ["claude", "-p", prompt],
            capture_output=True, text=True, timeout=120
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout.strip()
        return None
    except (FileNotFoundError, subprocess.TimeoutExpired) as e:
        print(f"  [Claude CLI failed: {e}]")
        return None


def _call_ollama(prompt: str) -> str | None:
    """Try Ollama local. Returns None on failure."""
    import urllib.request
    try:
        payload = json.dumps({
            "model": OLLAMA_MODEL,
            "prompt": prompt,
            "stream": False,
        }).encode()
        req = urllib.request.Request(
            f"{OLLAMA_URL}/api/generate",
            data=payload,
            headers={"Content-Type": "application/json"},
        )
        with urllib.request.urlopen(req, timeout=120) as resp:
            data = json.loads(resp.read())
            return data.get("response", "").strip() or None
    except Exception as e:
        print(f"  [Ollama failed: {e}]")
        return None


def call_llm(prompt: str) -> str:
    """Fallback chain: Gemini -> Claude -> Ollama."""
    global _backend_used

    result = _call_gemini(prompt)
    if result:
        _backend_used = "Gemini"
        return result

    print("  Gemini unavailable, trying Claude CLI...")
    result = _call_claude(prompt)
    if result:
        _backend_used = "Claude"
        return result

    print("  Claude unavailable, trying Ollama...")
    result = _call_ollama(prompt)
    if result:
        _backend_used = "Ollama"
        return result

    _backend_used = "none"
    return "[All LLM backends failed. Check GOOGLE_API_KEY, Claude CLI, or Ollama.]"


# ── Main processing ──────────────────────────────────────────────────────────

def process_folder(folder: Path):
    files = sorted([
        f for f in folder.iterdir()
        if f.is_file() and not f.name.startswith(".")
    ])

    if not files:
        print(f"No files found in {folder}")
        sys.exit(1)

    print(f"\nProcessing {len(files)} files from: {folder}")
    print(f"Output folder: {OUTPUT_DIR}\n")
    print("-" * 60)

    summaries = []

    for i, file_path in enumerate(files, 1):
        print(f"\n[{i}/{len(files)}] {file_path.name}")

        content, file_type = extract_content(file_path)
        char_count = len(content)
        print(f"  Type: {file_type} | Extracted: {char_count:,} chars")

        if "[Binary file" in content:
            print(f"  Skipping binary file")
            continue

        max_chars = 12000
        if len(content) > max_chars:
            content = content[:max_chars] + f"\n\n[... truncated — original was {len(content):,} chars]"

        prompt = ANALYSIS_PROMPT.format(
            filename=file_path.name,
            file_type=file_type,
            content=content
        )

        print(f"  Sending to LLM...")
        analysis = call_llm(prompt)
        print(f"  Backend: {_backend_used}")

        safe_name = file_path.stem.replace(" ", "_")
        out_path = OUTPUT_DIR / f"{safe_name}_summary.md"
        out_path.write_text(
            f"# {file_path.name}\n\n"
            f"**File type:** {file_type} | **Processed:** {TODAY} | **Via:** {_backend_used}\n\n"
            f"{analysis}\n"
        )
        print(f"  Saved: {out_path.name}")

        summaries.append({
            "filename": file_path.name,
            "file_type": file_type,
            "chars": char_count,
            "analysis": analysis,
            "backend": _backend_used,
        })

    print("\n" + "-" * 60)
    print("\nBuilding master summary...")

    master_parts = [
        f"# File Processing Report\n",
        f"**Source folder:** `{folder}`  ",
        f"**Processed:** {TODAY}  ",
        f"**Files analysed:** {len(summaries)}\n",
        f"---\n",
    ]

    for s in summaries:
        master_parts.append(f"## {s['filename']} `{s['file_type']}` via {s['backend']}\n")
        master_parts.append(s["analysis"])
        master_parts.append("\n---\n")

    master_path = OUTPUT_DIR / "MASTER_SUMMARY.md"
    master_path.write_text("\n".join(master_parts))
    print(f"Master summary: {master_path}")

    print(f"\nDone! Processed {len(summaries)} files.")
    print(f"   Output: {OUTPUT_DIR}")

    return OUTPUT_DIR


def main():
    if len(sys.argv) > 1:
        folder = Path(sys.argv[1])
    else:
        folder = BASE_DIR / "inbox"
        if not folder.exists():
            folder = BASE_DIR / "vault-template" / "inbox"

    if not folder.exists():
        print(f"Folder not found: {folder}")
        sys.exit(1)

    process_folder(folder)


if __name__ == "__main__":
    main()
